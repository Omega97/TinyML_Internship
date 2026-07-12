#!/usr/bin/env python3
"""Build SARDINE ICTP presentation (10 slides, Midnight Board theme)."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
PRESENTATIONS = PROJECT_ROOT / "presentations"
ASSETS = PRESENTATIONS / "assets"

# Midnight Board palette (slide bg matches SARDINE-logo-dark-small.png)
LOGO_GREY = RGBColor(0x1F, 0x1F, 0x1F)
SLIDE_BG = LOGO_GREY
CARD_BG = RGBColor(0x2A, 0x2A, 0x2A)
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TEXT = ICE
MUTED = RGBColor(0x94, 0xA3, 0xB8)

NAVY_HEX = "#1E2761"
ICE_HEX = "#CADCFC"
GOLD_HEX = "#D4A853"
LIGHT_HEX = "#F5F7FA"
MUTED_HEX = "#64748B"


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title: str, *, dark: bool = True) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY


def _add_subtitle(slide, text: str, y: float, *, dark: bool = False, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.color.rgb = MUTED


def _add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, *, size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = BODY_TEXT
        p.space_after = Pt(8)


def _add_speaker_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def _generate_acpl_chart(out_path: Path) -> None:
    data = {
        "NNUE pilot": json.loads((PROJECT_ROOT / "plots/nnue_d1_gate_acpl.json").read_text())["acpl"],
        "HCE": json.loads((PROJECT_ROOT / "plots/hce_d1_gate_acpl.json").read_text())["acpl"],
        "Sunfish": json.loads((PROJECT_ROOT / "plots/sunfish_d1_gate_acpl.json").read_text())["acpl"],
    }
    elos = {k: max(400, round(2855 - 10 * v)) for k, v in data.items()}

    labels = list(data.keys())
    acpls = [data[k] for k in labels]
    colors = [GOLD_HEX, MUTED_HEX, "#94A3B8"]

    fig, ax1 = plt.subplots(figsize=(8, 4.2), dpi=150)
    fig.patch.set_facecolor(LIGHT_HEX)
    ax1.set_facecolor(LIGHT_HEX)
    bars = ax1.bar(labels, acpls, color=colors, edgecolor="white", linewidth=1.2, width=0.55)
    ax1.set_ylabel("ACPL (centipawns)", fontsize=11, color=MUTED_HEX)
    ax1.set_title("Depth-1 gate · 16 self-play games · Stockfish 100 ms/move", fontsize=12, color=NAVY_HEX, pad=12)
    ax1.tick_params(colors=MUTED_HEX)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    ax1.grid(axis="y", color="#E2E8F0", linewidth=0.8)
    for bar, val in zip(bars, acpls):
        ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 25, f"{val:.0f}", ha="center", fontsize=10, color=NAVY_HEX)

    ax2 = ax1.twinx()
    ax2.plot(labels, [elos[k] for k in labels], color=NAVY_HEX, marker="o", linewidth=2, markersize=8)
    ax2.set_ylabel("Heuristic Elo", fontsize=11, color=NAVY_HEX)
    ax2.set_ylim(0, 1800)
    ax2.tick_params(colors=NAVY_HEX)
    ax2.spines["top"].set_visible(False)

    for i, k in enumerate(labels):
        ax2.annotate(f"~{elos[k]}", (i, elos[k]), textcoords="offset points", xytext=(0, 12), ha="center", fontsize=10, color=NAVY_HEX)

    plt.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, bbox_inches="tight", facecolor=LIGHT_HEX)
    plt.close(fig)


def _generate_accumulator_diagram(out_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(9, 3.8), dpi=150)
    fig.patch.set_facecolor(LIGHT_HEX)
    ax.set_facecolor(LIGHT_HEX)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 4)
    ax.axis("off")

    boxes = [
        (0.2, 2.2, 1.6, 1.0, "844 features\n(own POV)", ICE_HEX),
        (0.2, 0.6, 1.6, 1.0, "844 features\n(opp POV)", ICE_HEX),
        (2.4, 2.2, 1.5, 1.0, "Shared L1\n844 → W", NAVY_HEX),
        (2.4, 0.6, 1.5, 1.0, "Shared L1\n(same weights)", NAVY_HEX),
        (4.5, 1.4, 1.4, 1.2, "Concat\n2W", GOLD_HEX),
        (6.3, 1.4, 1.3, 1.2, "Router\n8 buckets", "#50808E"),
        (8.0, 1.4, 1.5, 1.2, "Expert\n2W → 1", NAVY_HEX),
    ]
    for x, y, w, h, label, color in boxes:
        rect = plt.Rectangle((x, y), w, h, facecolor=color, edgecolor="white", linewidth=2, zorder=2)
        ax.add_patch(rect)
        tc = "white" if color in (NAVY_HEX, "#50808E") else NAVY_HEX
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=9, color=tc, fontweight="bold")

    arrow_kw = dict(arrowstyle="->", color=NAVY_HEX, lw=1.8, mutation_scale=12)
    ax.annotate("", xy=(2.35, 2.7), xytext=(1.85, 2.7), arrowprops=arrow_kw)
    ax.annotate("", xy=(2.35, 1.1), xytext=(1.85, 1.1), arrowprops=arrow_kw)
    ax.annotate("", xy=(4.45, 2.0), xytext=(3.95, 2.7), arrowprops=arrow_kw)
    ax.annotate("", xy=(4.45, 1.8), xytext=(3.95, 1.1), arrowprops=arrow_kw)
    ax.annotate("", xy=(6.25, 2.0), xytext=(5.95, 2.0), arrowprops=arrow_kw)
    ax.annotate("", xy=(7.95, 2.0), xytext=(7.65, 2.0), arrowprops=arrow_kw)

    ax.text(5.0, 3.55, "Incremental add/sub on moves · full refresh on king centre-file crossing", ha="center", fontsize=10, color=MUTED_HEX)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, bbox_inches="tight", facecolor=LIGHT_HEX)
    plt.close(fig)


def _add_image(slide, path: Path, x: float, y: float, w: float, h: float | None = None) -> None:
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def _add_image_fit(slide, path: Path, x: float, y: float, max_w: float, max_h: float) -> None:
    """Place image centred in a box, preserving aspect ratio."""
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = iw / ih
    w, h = max_w, max_h
    if w / h > ratio:
        w = h * ratio
    else:
        h = w / ratio
    slide.shapes.add_picture(str(path), Inches(x + (max_w - w) / 2), Inches(y + (max_h - h) / 2), Inches(w), Inches(h))


def _add_table(slide, rows: list[list[str]], x: float, y: float, w: float, col_widths: list[float]) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(0.35 * len(rows))).table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(13 if ri > 0 else 14)
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else BODY_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else CARD_BG


def build() -> Path:
    PRESENTATIONS.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)

    acpl_chart = ASSETS / "acpl_d1_gate.png"
    acc_diagram = ASSETS / "accumulator_flow.png"
    _generate_acpl_chart(acpl_chart)
    _generate_accumulator_diagram(acc_diagram)

    logo = PROJECT_ROOT / "images" / "logo" / "SARDINE-logo-dark-small.png"
    arch = PROJECT_ROOT / "plots" / "sardine_nnue_architecture.png"
    sardine_img = ASSETS / (
        "sardine-fish-isolated-on-a-transparent-background-showcasing-its-shiny-scales-and-streamlined-body-"
        "a-sardine-fish-isolated-on-transparent-background-free-png.png"
    )
    anchovy_img = ASSETS / "anchovy-isolated-on-transparent-background-file-cut-out-ai-generated-png.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    s1 = prs.slides.add_slide(blank)
    _set_slide_bg(s1, SLIDE_BG)
    if logo.is_file():
        _add_image(s1, logo, 4.9, 0.85, 3.5)
    title_box = s1.shapes.add_textbox(Inches(0.55), Inches(3.95), Inches(12.2), Inches(0.85))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "SARDINE: Chess AI Under Hard Memory Limits"
    tp.font.name = "Trebuchet MS"
    tp.font.size = Pt(34)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    _add_subtitle(s1, "Small Artificial RAM-restricted Deep Intelligent Neural Engine", 4.85, dark=True, size=16)
    s1.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _add_subtitle(s1, "ICTP Internship · July 2026", 5.4, dark=True, size=14)
    s1.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _add_speaker_notes(s1, "Open with the constraint story: we target strong chess on a microcontroller with less RAM than a JPEG.")

    # --- Slide 2: Hardware ---
    s2 = prs.slides.add_slide(blank)
    _set_slide_bg(s2, SLIDE_BG)
    _add_title_bar(s2, "The Hardware Ceiling")
    _add_table(
        s2,
        [
            ["Resource", "Wio Terminal", "Why it matters"],
            ["CPU", "120 MHz (SAMD51)", "Eval latency caps search depth"],
            ["RAM", "192 KB", "TT must dominate; accumulators ~16 KB"],
            ["Flash", "512 KB", "Weights ~10% after int8 + L1 prune"],
            ["Move budget", "~1 second", "Urusov ESP32 baseline ~20 kNps (no NNUE)"],
            ["Target", "≥ 1700 Elo", "Dog (ESP32) proves feasibility at ~320 KB RAM"],
        ],
        0.55,
        1.25,
        12.0,
        [2.2, 3.5, 6.3],
    )
    stat = s2.shapes.add_shape(1, Inches(8.8), Inches(4.6), Inches(3.9), Inches(1.8))
    stat.fill.solid()
    stat.fill.fore_color.rgb = NAVY
    stat.line.fill.background()
    tf = stat.text_frame
    tf.paragraphs[0].text = "192 KB RAM"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "< one JPEG image"
    p2.font.size = Pt(14)
    p2.font.color.rgb = ICE
    p2.alignment = PP_ALIGN.CENTER
    _add_speaker_notes(s2, "Frame the problem viscerally: total RAM is smaller than a single photo file.")

    # --- Slide 3: Kaggle ---
    s3 = prs.slides.add_slide(blank)
    _set_slide_bg(s3, SLIDE_BG)
    _add_title_bar(s3, "What the Kaggle Challenge Taught Us")
    _add_subtitle(s3, "FIDE & Google Efficient Chess AI — 5 MiB RAM, 64 KiB binary cap", 1.05, dark=True)
    cards = [
        ("NNUE won at the top", "Micro int8 nets + king mirroring + geometric zeros beat pure HCE for peak Elo."),
        ("RAM = search depth", "Winners gutted history/hash tables to free 512 KB–1 MiB for the transposition table."),
        ("Buckets help", "2nd place Approvers: 8 output buckets by piece count in ~45 KB net."),
        ("C over C++", "Cfish base — lower runtime overhead; informs SARDINE's planned C port."),
    ]
    positions = [(0.55, 1.55), (6.75, 1.55), (0.55, 3.85), (6.75, 3.85)]
    for (title, body), (x, y) in zip(cards, positions):
        card = s3.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.9), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = NAVY
        tf = card.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.name = "Trebuchet MS"
        tf.paragraphs[0].font.size = Pt(17)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = GOLD
        pb = tf.add_paragraph()
        pb.text = body
        pb.font.name = "Calibri"
        pb.font.size = Pt(14)
        pb.font.color.rgb = BODY_TEXT
        pb.space_before = Pt(6)
    _add_speaker_notes(s3, "Cite linrock (1st), Approvers (2nd). Consensus: quantized NNUE + TT-heavy RAM + C implementation.")

    # --- Slide 4: Design principles ---
    s4 = prs.slides.add_slide(blank)
    _set_slide_bg(s4, SLIDE_BG)
    _add_title_bar(s4, "Design Principles We Adopted")
    _add_table(
        s4,
        [
            ["Principle", "Kaggle (5 MiB)", "SARDINE (192 KB)"],
            ["Eval", "Micro NNUE, int8, buckets", "844-dim bucketed NNUE + MoE (8 experts)"],
            ["RAM layout", "TT 512 KB – 1 MiB", "TT-dominant 128–160 KB"],
            ["Compression", "King mirror, pawn-rank zeros", "Same + gradual L1 prune 70–80%"],
            ["Runtime", "Cfish / C port", "C engine after PC bring-up"],
            ["Search", "LMR, null-move, SPSA", "Alpha-beta v0.3 → full stack + SPSA"],
        ],
        0.55,
        1.2,
        7.35,
        [1.9, 2.6, 2.85],
    )
    if sardine_img.is_file():
        _add_image_fit(s4, sardine_img, 8.35, 1.2, 4.45, 5.6)
    _add_speaker_notes(s4, "Same ideas as Kaggle winners, scaled down 25× on RAM.")

    # --- Slide 5: NNUE + MoE ---
    s5 = prs.slides.add_slide(blank)
    _set_slide_bg(s5, SLIDE_BG)
    _add_title_bar(s5, "Evaluation: Bucketed NNUE + MoE")
    if arch.is_file():
        arch_h = 4.35
        arch_w = arch_h * (1768 / 1330)
        arch_x = (13.333 - arch_w) / 2
        _add_image(s5, arch, arch_x, 1.15, arch_w, arch_h)
    _add_bullets(
        s5,
        [
            "844 sparse features (716 base + 128 tactical) · dual perspective",
            "Shared L1: 844 → W (W=128) · CReLU hidden · tanh LUT → expected reward [−1, +1]",
            "8 expert heads routed by piece count (+ queen-split at p ≤ 12)",
        ],
        0.55,
        5.65,
        12.0,
        1.35,
        size=14,
    )
    _add_speaker_notes(s5, "Teacher: Lc0 WDL (expected_reward = W − L). Production training via label_positions.py.")

    # --- Slide 6: Accumulator ---
    s6 = prs.slides.add_slide(blank)
    _set_slide_bg(s6, SLIDE_BG)
    _add_title_bar(s6, "Shared L1 Accumulator")
    _add_image(s6, acc_diagram, 0.55, 1.2, 12.0)
    _add_bullets(
        s6,
        [
            "Same pruned L1 weights called twice (own POV + mirrored opponent POV)",
            "Incremental add/sub on piece moves — bucket-agnostic shared layer",
            "Lazy full refresh only when the king crosses the centre file",
        ],
        0.55,
        5.5,
        12.0,
        1.5,
        size=14,
    )
    _add_speaker_notes(s6, "Incremental updates are why NNUE fits MCU budgets.")

    # --- Slide 7: Search ---
    s7 = prs.slides.add_slide(blank)
    _set_slide_bg(s7, SLIDE_BG)
    _add_title_bar(s7, "Alpha-Beta Search (v0.3)")
    col1 = s7.shapes.add_shape(1, Inches(0.55), Inches(1.2), Inches(5.8), Inches(5.5))
    col1.fill.solid()
    col1.fill.fore_color.rgb = CARD_BG
    col1.line.color.rgb = NAVY
    tf = col1.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Implemented now"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = GOLD
    now = [
        "Fixed-depth negamax alpha-beta",
        "Capture quiescence at leaf nodes",
        "MVV-LVA move ordering",
        "Perft validated: d5 = 4,865,609 nodes",
        "HCE default; NNUE via --eval nnue",
    ]
    for item in now:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_TEXT
        p.level = 0

    col2 = s7.shapes.add_shape(1, Inches(6.75), Inches(1.2), Inches(5.8), Inches(5.5))
    col2.fill.solid()
    col2.fill.fore_color.rgb = NAVY
    col2.line.fill.background()
    tf2 = col2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Roadmap"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = GOLD
    roadmap = [
        "Futility pruning + LMR + null-move",
        "Lazy evaluation (paired with accumulators)",
        "Iterative deepening with stable TT",
        "Killer moves (depth > 4)",
        "SPSA search-parameter tuning",
    ]
    for item in roadmap:
        p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = ICE
        p.level = 0
    _add_speaker_notes(s7, "Search skeleton is playable on PC; device port follows NNUE production train.")

    # --- Slide 8: Evaluation method ---
    s8 = prs.slides.add_slide(blank)
    _set_slide_bg(s8, SLIDE_BG)
    _add_title_bar(s8, "How We Estimate Strength: Stockfish ACPL Gate")
    steps = [
        ("1", "Engine self-play", "N games at fixed depth (currently 16)"),
        ("2", "Stockfish analysis", "100 ms/move — judges each played move"),
        ("3", "ACPL", "Average centipawn loss vs Stockfish best move"),
        ("4", "Heuristic Elo", "Elo ≈ 2855 − ACPL × 10 (floor 400, not FIDE)"),
        ("5", "Final gate", "cutechess-cli head-to-head for ≥ 1700 Elo (future)"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = 1.25 + i * 1.05
        circle = s8.shapes.add_shape(9, Inches(0.65), Inches(y), Inches(0.55), Inches(0.55))  # OVAL
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.paragraphs[0].text = num
        ctf.paragraphs[0].font.size = Pt(16)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

        box = s8.shapes.add_textbox(Inches(1.45), Inches(y - 0.05), Inches(6.55), Inches(0.85))
        btf = box.text_frame
        btf.word_wrap = True
        btf.paragraphs[0].text = title
        btf.paragraphs[0].font.name = "Trebuchet MS"
        btf.paragraphs[0].font.size = Pt(17)
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.color.rgb = WHITE
        bp = btf.add_paragraph()
        bp.text = desc
        bp.font.name = "Calibri"
        bp.font.size = Pt(13)
        bp.font.color.rgb = BODY_TEXT
    if anchovy_img.is_file():
        _add_image_fit(s8, anchovy_img, 8.15, 1.35, 4.65, 5.5)
    _add_speaker_notes(s8, "Method A1 from blueprint. Fast sanity check before expensive cutechess tournaments. Not win-rate vs Sunfish.")

    # --- Slide 9: Results ---
    s9 = prs.slides.add_slide(blank)
    _set_slide_bg(s9, SLIDE_BG)
    _add_title_bar(s9, "Early Results (16 self-play games, depth 1)")
    _add_image(s9, acpl_chart, 0.55, 1.15, 7.8)
    _add_bullets(
        s9,
        [
            "NNUE pilot ~1465 Elo (ACPL 139) — smoke net, not production",
            "HCE ~400 Elo (ACPL 357) — high variance, gross blunders",
            "Sunfish ~400 Elo (ACPL 1038) — calibration reference",
            "Pilot val_mse 0.056 · production Lichess+Lc0 train pending",
        ],
        8.6,
        1.4,
        4.1,
        4.5,
        size=13,
    )
    _add_speaker_notes(s9, "Few games → high variance. GIF demos available on request. NNUE d2 still weak until production net + search tuning.")

    # --- Slide 10: Status ---
    s10 = prs.slides.add_slide(blank)
    _set_slide_bg(s10, SLIDE_BG)
    _add_title_bar(s10, "Status & Next Steps", dark=True)
    items = [
        "✅  844-dim feature encoder (dual POV, 8 buckets, 94 tests)",
        "✅  Search v0.3 (alpha-beta, quiescence, MVV-LVA, perft)",
        "✅  NNUE pilot wired to engine (--eval nnue)",
        "✅  Stockfish ACPL gate + label_positions.py scaffold",
        "⬜  Lichess PGN → FEN + production nnue-pytorch training",
        "⬜  C port to Wio + incremental accumulators on device",
        "⬜  Full search stack + cutechess Elo gate (≥ 1700)",
    ]
    box = s10.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(8.5), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Calibri"
        p.font.size = Pt(17)
        p.font.color.rgb = ICE if item.startswith("⬜") else WHITE
        p.space_after = Pt(10)
    if logo.is_file():
        _add_image(s10, logo, 9.8, 4.8, 2.8)
    _add_subtitle(s10, "Thank you — questions?", 6.55, dark=True, size=20)
    _add_speaker_notes(s10, "Close with honest status: PC pipeline works; device deployment and Elo gate are next milestones.")

    out = PRESENTATIONS / "SARDINE_ICTP_2026-07.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        out = PRESENTATIONS / "SARDINE_ICTP_2026-07_rebuilt.pptx"
        prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")