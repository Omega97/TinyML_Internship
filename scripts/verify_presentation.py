#!/usr/bin/env python3
from pathlib import Path

from pptx import Presentation

p = Presentation(Path(__file__).resolve().parent.parent / "presentations" / "SARDINE_ICTP_2026-07.pptx")
for i, slide in enumerate(p.slides, 1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip().split("\n")[0][:70])
    has_notes = bool(slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip())
    title = texts[0] if texts else "(empty)"
    print(f"{i:2d}: {title}  {'[notes]' if has_notes else ''}")
print(f"Total slides: {len(p.slides)}")